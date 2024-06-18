import google.generativeai as genai
from pptx import Presentation


def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)

    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def generate_qna(text):
    genai.configure(api_key="AIzaSyDG4L7ze9_nZeGZ5SKFifBoDYkpiFTfPuk")
    response = genai.generate_text(
        model="models/text-bison-001",
        prompt="List 7 questions with answers based on the following:\n" + text,
        temperature=1,
        top_p=0.95,
        top_k=64,
        max_output_tokens=8192,
    )
    return response.result


if __name__ == "__main__":
    pptx_file = "OTP.pptx"

    extracted_text = extract_text_from_pptx(pptx_file)

    qna = generate_qna(extracted_text)

    print(qna)
