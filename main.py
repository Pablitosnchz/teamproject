import google.generativeai as genai
from pptx import Presentation

genai.configure(api_key="AIzaSyDG4L7ze9_nZeGZ5SKFifBoDYkpiFTfPuk") #Replace with your API key 

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)

    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    return text

def generate_notes(text):
  response = genai.generate_text(
      model="models/text-bison-001",  # Use text-bison for text generation 
      prompt="I want you to turn this text into well-written notes that are summarized and easily understandable:\n" + text,
      temperature=1,
      top_p=0.95,
      top_k=64,
      max_output_tokens=8192
  )
  return(response.result)

pptx_file = "OTP.pptx"

extracted_text = extract_text_from_pptx(pptx_file)

notes = generate_notes(extracted_text)

print(notes)